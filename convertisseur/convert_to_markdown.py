import os
import sys
import tkinter as tk
from tkinter import filedialog
import shutil

from docx import Document
from docx.shared import Inches
import pdfplumber
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re
import auto_labeling


def is_landscape(width, height):
    return width > height

def get_docx_orientation(docx_path):
    try:
        document = Document(docx_path)
        section = document.sections[0]
        width = section.page_width
        height = section.page_height
        return is_landscape(width, height)
    except Exception as e:
        print(f"Warning: Could not determine DOCX orientation: {e}")
        return False

def get_pdf_orientation(pdf_path):
    try:
        with pdfplumber.open(pdf_path) as pdf:
            first_page = pdf.pages[0]
            width = first_page.width
            height = first_page.height
            return is_landscape(width, height)
    except Exception as e:
        print(f"Warning: Could not determine PDF orientation: {e}")
        return False

def get_pptx_orientation(pptx_path):
    try:
        prs = Presentation(pptx_path)
        width = prs.slide_width
        height = prs.slide_height
        return is_landscape(width, height)
    except Exception as e:
        print(f"Warning: Could not determine PPTX orientation: {e}")
        return False

def remove_page_numbering(md_text):
    lines = md_text.splitlines()
    filtered_lines = [line for line in lines if not re.fullmatch(r'\s*\d+\s*', line)]
    return '\n'.join(filtered_lines)

def docx_table_to_markdown(table):
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
        rows.append(cells)
    md = []
    md.append('| ' + ' | '.join(rows[0]) + ' |')
    md.append('|' + '|'.join([' --- ' for _ in rows[0]]) + '|')
    for row in rows[1:]:
        md.append('| ' + ' | '.join(row) + ' |')
    return '\n'.join(md)

def extract_docx_images(docx_path, output_dir):
    image_links = []
    document = Document(docx_path)
    rels = document.part.rels
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    for rel in rels:
        rel = rels[rel]
        if "image" in rel.target_ref:
            image_part = rel.target_part
            image_name = os.path.basename(image_part.partname)
            image_path = os.path.join(output_dir, image_name)
            if not os.path.exists(image_path):
                with open(image_path, 'wb') as f:
                    f.write(image_part.blob)
            image_links.append(f"![{image_name}]({os.path.join(os.path.basename(output_dir), image_name)})")
    return image_links

def docx_to_markdown_text(docx_path, orientation_landscape, output_dir):
    document = Document(docx_path)
    markdown_lines = []
    header_label = "#slide" if orientation_landscape else "#page"
    markdown_lines.append(header_label)

    image_links = extract_docx_images(docx_path, output_dir)
    if image_links:
        markdown_lines.append('\n'.join(image_links))
        markdown_lines.append('')

    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style.name.startswith('Heading'):
            level = para.style.name.replace('Heading ', '')
            try:
                level_num = int(level)
            except:
                level_num = 1
            markdown_lines.append('#' * level_num + ' ' + text)
        else:
            markdown_lines.append(text)
        markdown_lines.append('')

    for table in document.tables:
        markdown_lines.append(docx_table_to_markdown(table))
        markdown_lines.append('')

    return '\n'.join(markdown_lines)

def pdf_table_to_markdown(table):
    md = []
    rows = table.extract_text().split('\n')
    for row in rows:
        cols = row.split()
        md.append('| ' + ' | '.join(cols) + ' |')
    return '\n'.join(md)

def pdf_to_markdown_text(pdf_path, orientation_landscape, output_dir):
    print(f"File is detected as a {'slide' if orientation_landscape else 'page'}")
    markdown_lines = []
    header_label = "#slide" if orientation_landscape else "#page"
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            markdown_lines.append(header_label + f" {i}")
            text = page.extract_text()
            if text:
                paragraphs = text.split('\n\n')
                for para in paragraphs:
                    para = para.strip()
                    if para:
                        markdown_lines.append(para)
                        markdown_lines.append('')
            tables = page.extract_tables()
            for table in tables:
                # Convert None to empty string in header row to avoid TypeError
                header_row = [cell if cell is not None else '' for cell in table[0]]
                md_table = '| ' + ' | '.join(header_row) + ' |\n'
                md_table += '|' + '|'.join([' --- ' for _ in header_row]) + '|\n'
                for row in table[1:]:
                    # Convert None to empty string to avoid TypeError
                    row_strs = [cell if cell is not None else '' for cell in row]
                    md_table += '| ' + ' | '.join(row_strs) + ' |\n'
                markdown_lines.append(md_table)
                markdown_lines.append('')
            # Extract images from the page
            images = page.images
            if images:
                if not os.path.exists(output_dir):
                    os.makedirs(output_dir)
                for idx, img in enumerate(images):
                    try:
                        # Extract image bytes
                        x0, top, x1, bottom = img['x0'], img['top'], img['x1'], img['bottom']
                        width = x1 - x0
                        height = bottom - top
                        # Crop image from page
                        cropped = page.within_bbox((x0, top, x1, bottom)).to_image(resolution=150)
                        image_path = os.path.join(output_dir, f"image_{i}_{idx}.png")
                        cropped.save(image_path, format="PNG")
                        markdown_lines.append(f"![image_{i}_{idx}]({os.path.join(os.path.basename(output_dir), f'image_{i}_{idx}.png')})")
                        markdown_lines.append('')
                    except Exception as e:
                        print(f"Warning: Could not extract image {idx} on page {i}: {e}")
    md_text = '\n'.join(markdown_lines)
    md_text = remove_page_numbering(md_text)
    return md_text
def extract_pptx_images(prs, output_dir):
    image_links = []
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    image_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                image_name = f"image_{image_count}.{image_ext}"
                image_path = os.path.join(output_dir, image_name)
                if not os.path.exists(image_path):
                    with open(image_path, 'wb') as f:
                        f.write(image_bytes)
                image_links.append(f"![{image_name}]({os.path.join(os.path.basename(output_dir), image_name)})")
                image_count += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                # Extract chart as image - workaround by exporting chart as image is not directly supported
                # So we skip chart extraction for now or add placeholder text
                image_name = f"chart_{image_count}.png"
                image_links.append(f"![{image_name}](#)")  # Placeholder link
                image_count += 1
    return image_links

def pptx_table_to_markdown(table):
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
        rows.append(cells)
    md = []
    md.append('| ' + ' | '.join(rows[0]) + ' |')
    md.append('|' + '|'.join([' --- ' for _ in rows[0]]) + '|')
    for row in rows[1:]:
        md.append('| ' + ' | '.join(row) + ' |')
    return '\n'.join(md)

def pptx_to_markdown_text(pptx_path, orientation_landscape, output_dir):
    prs = Presentation(pptx_path)
    markdown_lines = []
    header_label = "#slide" if orientation_landscape else "#page"
    markdown_lines.append(header_label)

    image_links = extract_pptx_images(prs, output_dir)
    if image_links:
        markdown_lines.append('\n'.join(image_links))
        markdown_lines.append('')

    slide_num = 1
    for slide in prs.slides:
        markdown_lines.append(f"{header_label} {slide_num}")
        slide_num += 1
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                markdown_lines.append(pptx_table_to_markdown(shape.table))
                markdown_lines.append('')
            elif hasattr(shape, "text") and shape.text.strip():
                markdown_lines.append(shape.text.strip())
                markdown_lines.append('')
    return '\n'.join(markdown_lines)

def extract_generic_images(input_file, output_dir):
    # Generic image extraction placeholder
    # For unsupported file types, try to extract images if possible
    # This is a stub and may require specific libraries per file type
    print(f"Generic image extraction not implemented for {input_file}")
    return []

def convert_to_markdown(input_file):
    ext = os.path.splitext(input_file)[1].lower()
    base_name = os.path.splitext(os.path.basename(input_file))[0]
    output_dir = os.path.join(os.path.dirname(input_file), base_name + "_images")
    if ext == '.docx':
        orientation_landscape = get_docx_orientation(input_file)
        markdown_content = docx_to_markdown_text(input_file, orientation_landscape, output_dir)
    elif ext == '.pdf':
        orientation_landscape = get_pdf_orientation(input_file)
        markdown_content = pdf_to_markdown_text(input_file, orientation_landscape, output_dir)
    elif ext == '.pptx':
        orientation_landscape = get_pptx_orientation(input_file)
        markdown_content = pptx_to_markdown_text(input_file, orientation_landscape, output_dir)
    else:
        # Attempt generic image extraction for other file types
        image_links = extract_generic_images(input_file, output_dir)
        markdown_content = ""
        if image_links:
            markdown_content += '\n'.join(image_links) + '\n\n'
        markdown_content += f"Unsupported file extension: {ext}. Only images extracted if any.\n"

    output_file = os.path.splitext(input_file)[0] + '.md'
    try:
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(markdown_content)
        print(f"Converted {ext} to Markdown successfully: {output_file}")

        # Call fine-tuned model script to enhance markdown with structure labeling
        labeled_output_file = os.path.splitext(input_file)[0] + '_labeled.md'
        try:
            auto_labeling.auto_label_markdown(output_file, labeled_output_file)
            print(f"Enhanced markdown with structure labeling: {labeled_output_file}")
        except Exception as e:
            print(f"Skipping auto-labeling due to error: {e}")

    except Exception as e:
        print(f"Error during conversion: {e}")

if __name__ == "__main__":
    root = tk.Tk()
    root.withdraw()
    filetypes = [("Supported files", "*.docx *.pdf *.pptx")]
    input_file = filedialog.askopenfilename(title="Select a file to convert to Markdown", filetypes=filetypes)
    if input_file:
        convert_to_markdown(input_file)
    else:
        print("No file selected. Exiting.")
